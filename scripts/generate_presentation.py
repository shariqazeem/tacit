#!/usr/bin/env python3
"""
ParallaxPay Hackathon Presentation Generator
Creates a beautiful PowerPoint presentation for x402 Solana Hackathon submission
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Solana brand colors
COLORS = {
    'primary': RGBColor(148, 0, 211),      # Purple
    'accent': RGBColor(0, 255, 163),       # Gradient Green
    'dark': RGBColor(20, 20, 35),          # Dark background
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(180, 180, 200),
    'highlight': RGBColor(138, 43, 226),   # Blue-Violet
}

def create_title_slide(prs):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    text_frame = title.text_frame
    text_frame.text = "ParallaxPay"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    text_frame = subtitle.text_frame
    text_frame.text = "Autonomous AI Agents on Distributed Compute\nwith x402 Micropayments"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['white']

    # Hackathon info
    info = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(0.8))
    text_frame = info.text_frame
    text_frame.text = "x402 Solana Hackathon 2024 • Parallax Eco Track"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['gray']

    return slide

def add_content_slide(prs, title, bullets, accent_color=None):
    """Add a content slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    text_frame = title_shape.text_frame
    text_frame.text = title

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = accent_color or COLORS['accent']

    # Add accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.4), Inches(2), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = accent_color or COLORS['accent']
    line.line.fill.background()

    # Content
    content = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.5), Inches(5))
    text_frame = content.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.space_before = Pt(12)
        p.space_after = Pt(12)

    return slide

def create_architecture_slide(prs):
    """Create architecture diagram slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    text_frame = title_shape.text_frame
    text_frame.text = "System Architecture"

    p = text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    # Architecture components
    components = [
        {"title": "AI Agents", "y": 2, "color": COLORS['highlight']},
        {"title": "x402 Protocol", "y": 3.2, "color": COLORS['accent']},
        {"title": "Parallax Cluster", "y": 4.4, "color": COLORS['primary']},
        {"title": "Solana Blockchain", "y": 5.6, "color": RGBColor(0, 204, 255)},
    ]

    for comp in components:
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(2), Inches(comp['y']), Inches(6), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = comp['color']
        shape.line.color.rgb = COLORS['white']
        shape.line.width = Pt(2)

        # Text
        text_frame = shape.text_frame
        text_frame.text = comp['title']
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']

        # Arrow (except for last component)
        if comp != components[-1]:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.5), Inches(comp['y'] + 0.85), Inches(1), Inches(0.25)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS['gray']
            arrow.line.fill.background()

    return slide

def create_demo_slide(prs):
    """Create demo/use cases slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']

    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    text_frame = title_shape.text_frame
    text_frame.text = "Live Demo Features"

    p = text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    # Feature boxes
    features = [
        {"title": "Market Oracle", "desc": "Autonomous crypto predictions\nwith multi-provider consensus", "x": 0.5},
        {"title": "Agent Swarms", "desc": "Collaborative intelligence\nfor provider discovery", "x": 3.5},
        {"title": "Micropayments", "desc": "Real-time x402 payments\non Solana", "x": 6.5},
    ]

    for feat in features:
        # Box
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(feat['x']), Inches(2.5), Inches(2.8), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary']
        shape.line.color.rgb = COLORS['accent']
        shape.line.width = Pt(3)

        # Title
        text_frame = shape.text_frame
        text_frame.text = feat['title']
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLORS['accent']
        p.space_after = Pt(10)

        # Description
        text_frame.add_paragraph()
        p = text_frame.paragraphs[1]
        p.text = feat['desc']
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['white']

        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Bottom stats
    stats = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(1))
    text_frame = stats.text_frame
    text_frame.text = "✓ 6 Agent Types  •  ✓ Multi-node Parallax Cluster  •  ✓ Real Solana Payments  •  ✓ MCP Server"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['gray']

    return slide

def create_closing_slide(prs):
    """Create closing slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['dark']

    # Title
    title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    text_frame = title.text_frame
    text_frame.text = "The Future of\nAutonomous AI Economies"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    # Subtitle
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
    text_frame = subtitle.text_frame
    text_frame.text = "ParallaxPay • Built on Gradient Parallax & Solana x402"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']

    # Call to action
    cta = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.8))
    text_frame = cta.text_frame
    text_frame.text = "github.com/shariqazeem/tacit"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['gray']

    return slide

def generate_presentation():
    """Generate the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("🎨 Creating title slide...")
    create_title_slide(prs)

    print("📊 Creating problem statement slide...")
    add_content_slide(prs, "The Problem", [
        "🤖 AI agents need compute resources but lack payment infrastructure",
        "💰 Traditional payment systems too slow and expensive for micropayments",
        "🌐 No decentralized marketplace connecting agents with compute providers",
        "🔒 Lack of trust and reputation system for autonomous transactions",
        "⚡ Existing solutions are centralized, manual, and don't scale"
    ])

    print("💡 Creating solution slide...")
    add_content_slide(prs, "ParallaxPay Solution", [
        "🎯 Autonomous AI Agent Marketplace on distributed compute infrastructure",
        "💸 x402 Protocol integration for seamless Solana micropayments ($0.001/inference)",
        "🏗️ Real distributed compute using Gradient Parallax multi-node clusters",
        "🤝 Swarm intelligence for collaborative provider discovery & consensus",
        "🏆 On-chain reputation system with trust badges and attestations",
        "⏰ Autonomous scheduling - agents run themselves without manual intervention"
    ])

    print("✨ Creating key features slide...")
    add_content_slide(prs, "Key Features", [
        "🔮 Market Oracle Agent: Autonomous crypto predictions with multi-provider consensus",
        "🐝 Agent Swarms: Multiple agents collaborate to benchmark and vote on providers",
        "⚙️ 6 Specialized Agents: Market Intel, Social Sentiment, DeFi Yield, Portfolio, Oracle, Blockchain Query",
        "📡 Real-time Provider Discovery: Auto-detect Parallax nodes with health checking",
        "💎 MCP Server: Model Context Protocol integration (bonus track!)",
        "🔄 Complete Autonomy: Self-scheduling, auto-payment, auto-execution"
    ], COLORS['primary'])

    print("🏗️ Creating architecture slide...")
    create_architecture_slide(prs)

    print("🎯 Creating Parallax integration slide...")
    add_content_slide(prs, "Gradient Parallax Integration", [
        "✅ Multi-node cluster deployment (1 scheduler + N workers)",
        "✅ Proper scheduler-worker architecture with load balancing",
        "✅ Real distributed inference (not simulated!)",
        "✅ Automatic failover and retry logic with health monitoring",
        "✅ OpenAI-compatible API with cost estimation",
        "✅ Comprehensive setup documentation and Docker deployment"
    ], COLORS['highlight'])

    print("💰 Creating x402 integration slide...")
    add_content_slide(prs, "x402 Protocol Integration", [
        "💸 Pay-per-inference model with automatic payment handling",
        "🔐 x402-fetch and x402-express middleware for seamless transactions",
        "📊 Real-time transaction tracking and public activity feed",
        "⛓️ On-chain verification with Solana Web3 integration",
        "📈 Complete payment history with wallet connection",
        "🚀 Production-ready payment infrastructure"
    ], COLORS['accent'])

    print("🎬 Creating demo slide...")
    create_demo_slide(prs)

    print("🏆 Creating competitive advantages slide...")
    add_content_slide(prs, "Why ParallaxPay Wins", [
        "🥇 Only submission with TRUE multi-node Parallax cluster (most use single instance)",
        "🤖 Real autonomous execution (self-scheduling agents, not manual clicks)",
        "🧠 Swarm intelligence with consensus algorithms (collaborative not competitive)",
        "🎯 Production-ready with Docker, SSL, comprehensive error handling",
        "🎁 Qualifies for multiple tracks: Parallax Eco + MCP Server",
        "📚 Exceptional documentation: 9 detailed docs, guides, troubleshooting",
        "💪 8,230+ lines of production TypeScript code with modern stack (Next.js 15, React 19)"
    ], RGBColor(255, 215, 0))

    print("✅ Creating closing slide...")
    create_closing_slide(prs)

    # Save presentation
    filename = "ParallaxPay_Hackathon_Presentation.pptx"
    prs.save(filename)
    print(f"\n✅ Presentation saved: {filename}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return filename

if __name__ == "__main__":
    print("🚀 ParallaxPay Presentation Generator")
    print("=" * 50)
    generate_presentation()
    print("\n🎉 Done! Ready for your hackathon submission!")
